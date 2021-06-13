import streamlit as st
from pptx import Presentation
from pptx.util import Inches, Pt
import os
import psutil
# import subprocess, sys
import webbrowser
import pandas as pd
import manipulate_news
import word_cloud
from datetime import datetime, timedelta
import seaborn as sns
import matplotlib.pyplot as plt
import re
import summarized_news

#Ref for slide types: 
# 0 ->  title and subtitle
# 1 ->  title and content
# 2 ->  section header
# 3 ->  two content
# 4 ->  Comparison
# 5 ->  Title only 
# 6 ->  Blank
# 7 ->  Content with caption
# 8 ->  Pic with caption

def ppt_insert_first_title(ppt_file, insert_title, insert_author, Layout=0, Placeholder=1, start_ppt=False):
    if os.path.exists(ppt_file):
        os.remove(ppt_file)
    prs=Presentation()

    title_slide_layout = prs.slide_layouts[Layout] #建立簡報檔第一張頁面物件
    #使用簡報物件中的方法將上一行建立的第一張頁面物件放進簡報
    slide = prs.slides.add_slide(title_slide_layout)
    #設定第一張頁面的標題 
    title = slide.shapes.title
    title.text = insert_title
    #設定第一張頁面的副標題
    subtitle = slide.placeholders[Placeholder] #設定副標題物件，副標題通常為第2個佔位圖
    subtitle.text = insert_author       

    #將簡報物件存檔
    prs.save(ppt_file)
    if start_ppt:
        # os.startfile(ppt_file)
        os.system("open "+ppt_file)

def ppt_insert_summarization(ppt_file, df, title_text, start_ppt=False):
    if os.path.exists(ppt_file):
        prs=Presentation(ppt_file)
    else:
        prs=Presentation()

    # title_slide_layout = prs.slide_layouts[Layout] #建立簡報檔第一張頁面物件
    #使用簡報物件中的方法將上一行建立的第一張頁面物件放進簡報
    #slide = prs.slides.add_slide(title_slide_layout)
    #shapes = slide.shapes

    # To create blank slide layout We have to use 6 as an argument of slide_layouts  
    for i in range(len(df)):
        if i % 3 == 0:
            # create one slide
            std_slide_layout = prs.slide_layouts[1] # title & content
            slide = prs.slides.add_slide(std_slide_layout)
            title = slide.shapes.title
            title.text = title_text
            title.text_frame.paragraphs[0].font.size = Pt(40)
            # add content
            body_shape = slide.shapes.placeholders[1]
            tf = body_shape.text_frame
            # header
            tf.text = df.loc[i, 'header'] 
            tf.paragraphs[0].font.size = Pt(18)
            tf.paragraphs[0].font.bold = True
        else:
            p = tf.add_paragraph()            
            # header
            p.text = df.loc[i, 'header']
            p.font.bold = True
            p.font.size = Pt(18)
            # p.level = 0

        # time & source & sentiment
        p = tf.add_paragraph()
        if df.loc[i, 'sentiment'] == 1:
            sentiment = 'POSITIVE'
        elif df.loc[i, 'sentiment'] == -1:
            sentiment = 'NEGATIVE'
        else:
            sentiment = 'NEUTRAL'
        p.text = df.loc[i, 'time_ago'] + ' | '+ df.loc[i, 'source'] + ' | ' + sentiment
        p.font.size = Pt(12)
        p.font.italic = True
        p.level = 1
        # content
        p = tf.add_paragraph()
        p.text = df.loc[i, 'content_summary']
        p.font.size = Pt(14)
        p.level = 1

    #將簡報物件存檔
    prs.save(ppt_file)
    if start_ppt:
        # os.startfile(ppt_file)
        os.system("open "+ppt_file)

def ppt_insert_images(ppt_file, img_title, img_path, start_ppt=False, height = Inches(6.5), left = Inches(0.73), top = Inches(1.2)):
    if os.path.exists(ppt_file):
        prs=Presentation(ppt_file)
    else:
        prs=Presentation()

    slide = prs.slides.add_slide(prs.slide_layouts[5]) # title only
    shapes = slide.shapes
    #投影片標題
    title_shape = shapes.title
    title_shape.text = img_title # image_profiles[0][0]    
    title_shape.text_frame.paragraphs[0].font.size = Pt(40)

    # show the figure
    pic = slide.shapes.add_picture(img_path, left, top, height=height)

    #將簡報物件存檔
    prs.save(ppt_file)
    if start_ppt:
        # os.startfile(ppt_file)
        # opener = "open" if sys.platform == "darwin" else "xdg-open"
        # subprocess.call([opener, ppt_file])
        os.system("open "+ppt_file)

def ppt_insert_sentiment(ppt_file, title_text, postive5, negative5, start_ppt=False):
    if os.path.exists(ppt_file):
        prs=Presentation(ppt_file)
    else:
        prs=Presentation()


    slide = prs.slides.add_slide(prs.slide_layouts[3]) # 2 contents
    #投影片標題
    title = slide.shapes.title
    title.text = title_text
    title.text_frame.paragraphs[0].font.size = Pt(40)

    # positive 5 companies
    body_shape1 = slide.shapes.placeholders[1]
    tf1 = body_shape1.text_frame   
    tf1.text = 'Highest Positive Rate'
    tf1.paragraphs[0].font.size = Pt(24)
    tf1.paragraphs[0].font.bold = True
    for i in range(5):
        p = tf1.add_paragraph()
        p.text = postive5[i]
        p.font.size = Pt(18)
        p.level = 1
    # negative 5 companies
    body_shape2 = slide.shapes.placeholders[2]
    tf2 = body_shape2.text_frame   
    tf2.text = 'Highest Negative Rate'
    tf2.paragraphs[0].font.size = Pt(26)
    tf2.paragraphs[0].font.bold = True
    for i in range(5):
        p = tf2.add_paragraph()
        p.text = negative5[i]
        p.font.size = Pt(20)
        p.level = 1
    
    #將簡報物件存檔
    prs.save(ppt_file)
    if start_ppt:
        # os.startfile(ppt_file)
        # opener = "open" if sys.platform == "darwin" else "xdg-open"
        # subprocess.call([opener, ppt_file])
        os.system("open "+ppt_file)

def app():
    st.title('Automated PPT Generation')

    for proc in psutil.process_iter():
        if proc.name() == 'POWERPNT.EXE':
            proc.kill()
    
    sp500 = pd.read_csv('data/constituents_csv.csv')
    sp500.loc[:, 'name_clean'] = [re.sub('\s((Brands\s)?Inc\.?|Company|Corp\.?|Bancorp|Technologies|\&?\s?Co\.|Entertainment|Corporation|Svc\.Gp\.)$', '', n) for n in sp500.Name]
    col1, _ = st.beta_columns((2,1))
    title = col1.text_input("Title: (optional)", "S&P500 News Summary")
    file_name = col1.text_input("File Name: (optional)", "news_summary_t2")
    latest_news = col1.checkbox("Latest Unbiased News")
    sentiment_analysis = col1.checkbox("Sentiment Analysis")
    esg = col1.checkbox("ESG Analysis")

    col1, col2 = st.beta_columns((2,1))
    select_category = col1.multiselect('Select categories: (optional)', sp500.Sector.unique().tolist(), [])
    category_n_news = col2.number_input('# news', min_value = 0, max_value = 12, value = 0, step = 1, key = 'cat')
    select_companies = col1.multiselect('Select companies: (optional)', sp500.name_clean.tolist(), [])
    companies_n_news = col2.number_input('# news', min_value = 0, max_value = 12, value = 0, step = 1, key = 'comp')
    

    go = st.button('Generate')

    if go:
        # read data
        data_news = pd.read_json('data/data_bias_news.json')
        top6_news_df = manipulate_news.select_news(data_news, n = 6)
        data_news['time'] = pd.to_datetime(data_news['time'], unit='ms')  
        data_news['date'] = data_news['time'].dt.date

        #### page 1 ####
        ppt_insert_first_title(ppt_file=file_name+'.pptx', insert_title=title, insert_author='Team 2', start_ppt=False)
        
        if latest_news:
            #### page 2: wordcloud ####
            wc_fig = word_cloud.plot_wordcloud(data_news, n_words = 100, date = datetime.today(), set_width = 1200, set_height = 800)
            wc_fig.savefig('report/img/wc_fig.png')
            ppt_insert_images(ppt_file=file_name+'.pptx', 
                            img_title = 'Key words of this week', 
                            img_path = 'report/img/wc_fig.png', 
                            start_ppt=False)

            #### page 3-5: summarization ####
            ppt_insert_summarization(ppt_file=file_name+'.pptx', 
                                    df = top6_news_df,
                                    title_text = 'Latest Unbiased News',
                                    start_ppt=False)
        if sentiment_analysis:
            #### page 6: sentiment -- company ####
            df_positive = pd.read_json('data/data_entities_pos_rate.json')
            # sort df_positive data by total news count (50 companies which have most news)
            df_positive_top50_sum = df_positive.sort_values(by = 'sum', ascending = False).reset_index(drop = True)[:50]
            # top5 positive companies
            top5_positive = df_positive_top50_sum.sort_values(by = 'positive_rate', ascending=False)[:5].reset_index(drop = True)
            # top5 negative companies
            df_else = df_positive_top50_sum.sort_values(by = 'positive_rate', ascending=False)[5:]
            top5_negative = df_else.sort_values(by = 'negative_rate', ascending=False)[:5].reset_index(drop = True)
            ppt_insert_sentiment(ppt_file=file_name+'.pptx',
                                title_text = 'Sentiment Analysis -- Company', 
                                postive5 = top5_positive['entities'].tolist(), 
                                negative5 = top5_negative['entities'].tolist(), 
                                start_ppt=False)

            #### page 7: sentiment -- source ####
            # count sentiment of source per day
            source_sen_count = data_news.groupby(by = ['date', 'source', 'sentiment']).count()[['id']].reset_index()
            source_sen_count_p = source_sen_count.pivot(index = ['date', 'source'], columns = 'sentiment', values = 'id').reset_index()
            source_sen_count_p.fillna(0, inplace = True)
            # total news count (in that day)
            source_sen_count_p['sum'] = source_sen_count_p[[-1, 0, 1]].sum(axis = 1)
            # score of the day for every source
            source_sen_count_p['score'] = (source_sen_count_p[1] - source_sen_count_p[-1]) / source_sen_count_p['sum']
            # select top 10 source (which have the most # news)
            top10_source = data_news.groupby('source').count()[['id']].sort_values(by='id', ascending=False)[:10].reset_index()['source'].tolist()
            top10_sen_count = source_sen_count_p[source_sen_count_p.source.isin(top10_source)]
            ### box plot of score
            plt.figure(figsize=(8,6))
            sns.boxplot(data = top10_sen_count, x = 'source', y = 'score', palette='Set3')
            plt.xticks(rotation = -45)
            plt.tight_layout()
            plt.savefig('report/img/source_sentiment_score.png')
            
            ppt_insert_images(ppt_file=file_name+'.pptx', 
                            img_title = 'Sentiment Analysis -- Source', 
                            img_path = 'report/img/source_sentiment_score.png', 
                            start_ppt=False, 
                            height = Inches(5.5),
                            left = Inches(1.25),
                            top = Inches(1.5))


        #### page 8: ESG ####
        if esg:
            pass
        
        #### page 9: selected sectors ####
        if select_category:            
            df_week = data_news[data_news.time >= datetime.today() - timedelta(days=7)]
            # manipulate company name (explode)
            df_week['company'] = df_week.company_all.copy()
            df_exploded = df_week.explode('company')
            df_exploded.dropna(subset = ['company'], inplace = True)
            # replace Symbol & full name with clean name
            df_exploded.company_all.replace(sp500.Symbol.tolist(), sp500.name_clean.tolist(), inplace = True)
            df_exploded.company_all.replace(sp500.Name.tolist(), sp500.name_clean.tolist(), inplace = True)
            # merge with sector
            df_exploded = df_exploded.merge(sp500, how = 'left', left_on = 'company', right_on = 'name_clean')
            # all selected categories news
            display_df = df_exploded.loc[df_exploded.Sector.isin(select_category)]
            display_df.drop_duplicates(subset = ['header'], inplace = True)
            
            ## wordcloud
            cat_wc_fig = word_cloud.plot_wordcloud(display_df, n_words = 100, date = None, set_width = 1200, set_height = 800)
            cat_wc_fig.savefig('report/img/cat_wc_fig.png')
            ppt_insert_images(ppt_file=file_name+'.pptx', 
                            img_title = 'Key words of the selected categories', 
                            img_path = 'report/img/cat_wc_fig.png', 
                            start_ppt=False)

            ## summarization 
            summarized_df = summarized_news.summarized_multiple_news(display_df, n_sen = category_n_news)
            summarized_df_time = manipulate_news.calculate_time(summarized_df)
            ppt_insert_summarization(ppt_file=file_name+'.pptx', 
                                    df = summarized_df_time,
                                    title_text = 'Selected Categories News',
                                    start_ppt=False)

        #### page 10: selected companies ####
        if select_companies:
            company_news_df = pd.DataFrame()
            for i in range(len(select_companies)):
                company_i = sp500[sp500.name_clean == select_companies[i]].index[0]
                company_news_df = pd.concat([company_news_df, data_news[data_news.company_all.apply(lambda x: (select_companies[i] in x) or (sp500.loc[company_i, 'Symbol'] in x) or (sp500.loc[company_i, 'Name'] in x))]])
            ## wordcloud
            company_wc_fig = word_cloud.plot_wordcloud(company_news_df, n_words = 100, date = None, set_width = 1200, set_height = 800)
            company_wc_fig.savefig('report/img/company_wc_fig.png')
            ppt_insert_images(ppt_file=file_name+'.pptx', 
                            img_title = 'Key words of the selected companies', 
                            img_path = 'report/img/company_wc_fig.png', 
                            start_ppt=False)
            ## summarization 
            summarized_df_com = summarized_news.summarized_multiple_news(company_news_df, n_sen = companies_n_news)
            summarized_df_time_com = manipulate_news.calculate_time(summarized_df_com)
            ppt_insert_summarization(ppt_file=file_name+'.pptx', 
                                    df = summarized_df_time_com,
                                    title_text = 'Selected Companies News',
                                    start_ppt=False)


        os.system("open "+file_name+'.pptx')
